# backend/app/agents/experience.py
import io                                         # 바이트 스트림 처리
from app.core.gemini import call_gemini                          # Gemini API
from app.core.config import settings              # 환경변수




def extract_text_from_pdf(content: bytes) -> str:
    """PDF에서 텍스트 추출"""
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() for page in reader.pages)

def extract_text_from_docx(content: bytes) -> str:
    """Word 파일에서 텍스트 추출"""
    import docx
    doc = docx.Document(io.BytesIO(content))
    return "\n".join(para.text for para in doc.paragraphs)

def extract_text_from_pptx(content: bytes) -> str:
    """PPT 파일에서 텍스트 추출"""
    from pptx import Presentation                 # PPT 파싱 라이브러리
    prs = Presentation(io.BytesIO(content))
    texts = []
    for slide in prs.slides:                      # 슬라이드 순회
        for shape in slide.shapes:                # 슬라이드 내 도형 순회
            if shape.has_text_frame:              # 텍스트가 있는 도형만
                texts.append(shape.text_frame.text)
    return "\n".join(texts)

def extract_text_by_type(content: bytes, content_type: str, filename: str) -> str:
    """파일 타입에 따라 자동 분기"""
    filename_lower = filename.lower()
    if filename_lower.endswith(".pdf") or "pdf" in content_type:
        return extract_text_from_pdf(content)
    elif filename_lower.endswith(".docx") or "word" in content_type:
        return extract_text_from_docx(content)
    elif filename_lower.endswith(".pptx") or "presentation" in content_type:
        return extract_text_from_pptx(content)
    else:
        return content.decode("utf-8", errors="ignore")  # 그 외는 텍스트로 처리

async def experience_agent(
    raw_text: str = None,      # 직접 입력한 경험 텍스트
    file_text: str = None,     # 파일에서 추출한 텍스트
) -> dict:
    """경험 분석 에이전트 — 입력된 경험에서 핵심 키워드 자동 추출"""

    # 텍스트 소스 합치기 (둘 다 있으면 합산)
    combined = ""
    if raw_text:
        combined += f"[직접 입력한 경험]\n{raw_text}\n\n"
    if file_text:
        combined += f"[첨부 파일 내용]\n{file_text[:5000]}\n\n"  # 5000자 제한

    if not combined:
        return {"keywords": [], "summary": ""}

    prompt = f"""
너는 채용 컨설턴트야. 아래는 지원자의 경험 정보야.
이 내용을 분석해서 AI/SW 직무 지원에 활용할 수 있는 핵심 경험 키워드를 추출해줘.

{combined}

다음 항목을 JSON 형식으로 작성해줘:
1. keywords: 핵심 경험 키워드 리스트 (10개 이내, 기술/역량/프로젝트 경험 포함)
2. summary: 지원자 경험 한 줄 요약
3. strengths: 강점 3가지 (자소서에 바로 쓸 수 있는 수준으로)

JSON만 반환하고 다른 말은 하지 마.
"""

    from app.core.gemini import call_gemini, parse_json_response

    text = call_gemini(prompt)
    parsed = parse_json_response(text)

    return {
        "agent": "experience",
        "keywords": parsed.get("keywords", []),    # 추출된 키워드
        "summary": parsed.get("summary", ""),      # 경험 요약
        "strengths": parsed.get("strengths", [])   # 강점
    }